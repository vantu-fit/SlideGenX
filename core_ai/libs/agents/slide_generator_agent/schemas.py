import os, re
from typing import List, Dict, Any, Optional, Union
from pydantic import BaseModel, Field, ConfigDict
from pptx import Presentation
import json
import logging


logger = logging.getLogger(__name__)


class PlaceHolder(BaseModel):
    placeholder: Optional[int] = Field(description="Placeholder index")
    content: Optional[Union[str, List[str]]] = Field(description="Content for the placeholder")

class Content(BaseModel):
    slide: int = Field(description="Slide number")
    slide_content: Optional[List[PlaceHolder]] = Field(description="List Content for the slide")
    

class RefitContent(BaseModel):
    contents: List[Content] = Field(description="List of content items to be refitted")

class Map(BaseModel):
    mappings: Dict[str, Union[str, List[str]]] = Field()


class SlideMapping(BaseModel):
    """Model for mapping slide content to placeholders"""
    slides : List[Map] = Field(description="Mapping of slide content to placeholders content")

class ChooseLayout(BaseModel):
    """ ChooseLayout class to hold layout information """
    slide_indexs : List[int] = Field(description="List of slide indices")


class SlideContent(BaseModel):
    """Model for individual slide content"""
    model_config = ConfigDict(title="SlideContent", 
                             extra="allow",
                             arbitrary_types_allowed=True)
    
    title: Optional[str] = Field(default=None, description="Slide title")
    layout_index: int = Field(description="Index of the layout to use (1-based)")
    layout_name: str = Field(description="Name of the layout")
    content: Dict[str, Any] = Field(description="Content for each placeholder")
    notes: Optional[str] = Field(default=None, description="Speaker notes for this slide")

class PresentationContent(BaseModel):
    """Model for the full presentation content"""
    model_config = ConfigDict(title="PresentationContent", 
                             extra="allow",
                             arbitrary_types_allowed=True)
    
    slides: List[SlideContent] = Field(description="List of slides to create")
    title: Optional[str] = Field(default=None, description="Presentation title")
    theme: Optional[Dict[str, Any]] = Field(default=None, description="Theme information")
    
    @classmethod
    def from_json(cls, json_data: Dict[str, Any]) -> "PresentationContent":
        """Create PresentationContent instance from JSON dictionary"""
        return cls.model_validate(json_data)
    
    @classmethod
    def from_json_file(cls, file_path: str) -> "PresentationContent":
        """Create PresentationContent instance from JSON file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return cls.from_json(data)
    
    def to_json(self) -> Dict[str, Any]:
        """Convert to JSON dictionary"""
        return self.model_dump(exclude_none=True)
    
    def to_json_file(self, file_path: str) -> None:
        """Save to JSON file"""
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(self.to_json(), f, ensure_ascii=False, indent=2)

class SlideLayoutManager:
    """Class to manage slide layouts from a PowerPoint template"""
   
    def __init__(self, template_path: str):
        """Initialize with a PowerPoint template path"""
        self.template_path = template_path
        self.layouts = self._load_layouts()
   
    def _load_layouts(self) -> Dict[int, Dict]:
        """Load layout information from the template"""
        prs = Presentation(self.template_path)
        layouts = {}
       
        for i, layout in enumerate(prs.slide_layouts):
            layout_info = {
                'index': i + 1,
                'name': layout.name,
                'placeholders': []
            }
           
            for ph in layout.placeholders:
                ph_info = {
                    'idx': ph.placeholder_format.idx,
                    'type': ph.placeholder_format.type,
                    'name': ph.name
                }
                layout_info['placeholders'].append(ph_info)
           
            layouts[i + 1] = layout_info
           
        return layouts
   
    def get_layout_info(self) -> str:
        """Get a formatted string of all layouts for prompting"""
        result = f"INFORMATION ABOUT LAYOUTS IN TEMPLATE {os.path.basename(self.template_path)}:\n\n"
       
        for idx, layout in self.layouts.items():
            result += f"Layout {idx} - {layout['name']} :\n"
            result += f"  Contains {len(layout['placeholders'])} placeholder(s)\n"
           
            for ph in layout['placeholders']:
                result += f"  * {ph['name']} (Type: {ph['type']})\n"
           
            result += "\n"
           
        return result
    
    # This would be added to the SlideLayoutManager class
    def get_placeholder_map(self, layout_index: int) -> Dict[str, int]:
        """
        Get the placeholder mapping for a specific layout
        
        Args:
            layout_index: The 1-based index of the layout
            
        Returns:
            Dictionary mapping placeholder names to types
        """
        layout_map = self.create_map()
        
        # Layout indices in the map are 0-based, but input is 1-based
        adjusted_index = layout_index - 1
        
        if adjusted_index < 0 or adjusted_index >= len(layout_map):
            logger.warning(f"Invalid layout index {layout_index}, defaulting to first layout")
            return layout_map[0]
        
        return layout_map[adjusted_index]

    def create_map(self, indexs: List[int] = None):
        """
        Create a mapping of layouts and their placeholders
        
        Args:
            indexs: Optional list of layout indices to include (1-based). If provided,
                only return the mapping for these specific layouts.
        
        Returns:
            List of dictionaries mapping placeholder names to types
        """
        map_list = []

        indexs = indexs.model_dump().get('slide_indexs') if not isinstance(indexs, list) else indexs

        
        # Get all layouts sorted by index
        sorted_layouts = sorted(self.layouts.items(), key=lambda x: x[0])
        # If specific indices are provided, filter the layouts
        if indexs and len(indexs) > 0:
            sorted_layouts = [(idx, layout) for idx, layout in sorted_layouts if idx in indexs]
        
        # Create the mapping for each layout
        for idx, layout in sorted_layouts:
            item = {}
            # Create a mapping of placeholder names to types
            for ph in layout['placeholders']:
                key = ph['name']
                val = ph['type']  # Already an integer, no need to convert
                item[key] = val
            
            # Add this layout's mapping to the result list
            map_list.append(item)
        
        return map_list
   
    def get_layouts_by_name(self, name: str) -> List[int]:
        """Find layouts by name (case-insensitive partial match)"""
        matches = []
        pattern = re.compile(name.lower())
       
        for idx, layout in self.layouts.items():
            if pattern.search(layout['name'].lower()):
                matches.append(idx)
               
        return matches

if __name__ == "__main__":
    layout_manager = SlideLayoutManager("pptx_templates/Ion_Boardroom.pptx")
    layout_info = layout_manager.get_layout_info()
    print(layout_info)
    print(json.dumps(layout_manager.create_map([10,1]), indent=2))