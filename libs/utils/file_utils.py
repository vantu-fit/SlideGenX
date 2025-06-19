"""
Utility functions for file operations.
"""

import os
import json
import logging
from typing import Dict, Any, Optional, Union, List

logger = logging.getLogger(__name__)


def ensure_directory_exists(directory_path: str) -> bool:
    """
    Ensure that a directory exists, creating it if necessary.
    
    Args:
        directory_path: Path to the directory
        
    Returns:
        True if the directory exists or was created, False otherwise
    """
    try:
        if not os.path.exists(directory_path):
            os.makedirs(directory_path)
            logger.info(f"Created directory: {directory_path}")
        return True
    except Exception as e:
        logger.error(f"Failed to create directory {directory_path}: {e}")
        return False


def save_json(data: Dict[str, Any], file_path: str, ensure_dir: bool = True) -> bool:
    """
    Save data as JSON to a file.
    
    Args:
        data: Data to save
        file_path: Path where to save the data
        ensure_dir: Whether to ensure the directory exists
        
    Returns:
        True if successful, False otherwise
    """
    try:
        # Ensure directory exists if requested
        if ensure_dir:
            directory = os.path.dirname(file_path)
            if directory and not ensure_directory_exists(directory):
                return False
        
        # Save data
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        
        logger.info(f"Saved JSON data to {file_path}")
        return True
    except Exception as e:
        logger.error(f"Failed to save JSON data to {file_path}: {e}")
        return False


def load_json(file_path: str) -> Optional[Dict[str, Any]]:
    """
    Load JSON data from a file.
    
    Args:
        file_path: Path to the JSON file
        
    Returns:
        Loaded data, or None if loading failed
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        logger.info(f"Loaded JSON data from {file_path}")
        return data
    except Exception as e:
        logger.error(f"Failed to load JSON data from {file_path}: {e}")
        return None


def export_presentation_to_pptx(presentation_data: Dict[str, Any], output_path: str) -> bool:
    """
    Export presentation data to a PowerPoint file.
    
    Args:
        presentation_data: The presentation data
        output_path: Path where to save the PowerPoint file
        
    Returns:
        True if successful, False otherwise
    """
    try:
        # Import python-pptx
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # Create a new presentation
        prs = Presentation()
        
        # Get presentation outline and style
        outline = presentation_data.get("outline", {})
        style = presentation_data.get("style", {})
        slides = presentation_data.get("slides", [])
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        # Set title
        title = title_slide.shapes.title
        title.text = outline.get("title", "Presentation")
        
        # Set subtitle
        subtitle = title_slide.placeholders[1]
        subtitle.text = outline.get("overall_message", "")
        
        # Add content slides
        for slide_data in slides:
            # Skip the title slide (already added)
            if slide_data.get("slide_content", {}).get("slide_number", 0) == 1:
                continue
                
            # Add a new slide
            content_slide_layout = prs.slide_layouts[1]  # Layout with title and content
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("slide_content", {}).get("title", "")
            
            # Set content
            content = slide.placeholders[1]
            content.text = slide_data.get("slide_content", {}).get("content", "")
        
        # Save the presentation
        prs.save(output_path)
        
        logger.info(f"Exported presentation to PowerPoint: {output_path}")
        return True
    except ImportError:
        logger.error("python-pptx not installed. Cannot export to PowerPoint.")
        return False
    except Exception as e:
        logger.error(f"Failed to export presentation to PowerPoint: {e}")
        return False


def export_presentation_to_html(presentation_data: Dict[str, Any], output_path: str) -> bool:
    """
    Export presentation data to an HTML file.
    
    Args:
        presentation_data: The presentation data
        output_path: Path where to save the HTML file
        
    Returns:
        True if successful, False otherwise
    """
    try:
        # Get presentation outline and style
        outline = presentation_data.get("outline", {})
        style = presentation_data.get("style", {})
        slides = presentation_data.get("slides", [])
        
        # Create HTML content
        html_content = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{outline.get("title", "Presentation")}</title>
            <style>
                body {{
                    font-family: {style.get("font_family", "Arial, sans-serif")};
                    margin: 0;
                    padding: 0;
                    background-color: #f5f5f5;
                }}
                .slide {{
                    width: 800px;
                    height: 600px;
                    margin: 20px auto;
                    padding: 40px;
                    background-color: white;
                    box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
                    position: relative;
                    overflow: hidden;
                }}
                .slide-title {{
                    font-size: 36px;
                    margin-bottom: 20px;
                    color: #333;
                }}
                .slide-content {{
                    font-size: 24px;
                    line-height: 1.5;
                }}
                .title-slide {{
                    display: flex;
                    flex-direction: column;
                    justify-content: center;
                    align-items: center;
                    text-align: center;
                }}
                .title-slide .slide-title {{
                    font-size: 48px;
                    margin-bottom: 40px;
                }}
                .notes {{
                    margin-top: 20px;
                    padding: 10px;
                    background-color: #f9f9f9;
                    border-left: 4px solid #ddd;
                    font-size: 16px;
                }}
            </style>
        </head>
        <body>
        """
        
        # Add slides
        for slide_data in slides:
            slide_content = slide_data.get("slide_content", {})
            slide_number = slide_content.get("slide_number", 0)
            slide_title = slide_content.get("title", "")
            slide_content_text = slide_content.get("content", "")
            slide_notes = slide_content.get("notes", "")
            
            # Determine if this is the title slide
            is_title_slide = (slide_number == 1)
            
            # Create slide HTML
            html_content += f"""
            <div class="slide {'title-slide' if is_title_slide else ''}">
                <h1 class="slide-title">{slide_title}</h1>
                <div class="slide-content">{slide_content_text.replace('\n', '<br>')}</div>
                
                {f'<div class="notes">Speaker Notes: {slide_notes}</div>' if slide_notes else ''}
            </div>
            """
        
        # Close HTML
        html_content += """
        </body>
        </html>
        """
        
        # Save the HTML file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        logger.info(f"Exported presentation to HTML: {output_path}")
        return True
    except Exception as e:
        logger.error(f"Failed to export presentation to HTML: {e}")
        return False


def export_presentation(presentation_data: Dict[str, Any], output_path: str, format: str = "auto") -> bool:
    """
    Export presentation data to the specified format.
    
    Args:
        presentation_data: The presentation data
        output_path: Path where to save the exported file
        format: Export format ("pptx", "html", "json", or "auto" to determine from extension)
        
    Returns:
        True if successful, False otherwise
    """
    # Determine format from file extension if set to auto
    if format == "auto":
        _, extension = os.path.splitext(output_path)
        format = extension.strip(".").lower()
    
    # Export based on format
    if format == "pptx":
        return export_presentation_to_pptx(presentation_data, output_path)
    elif format == "html":
        return export_presentation_to_html(presentation_data, output_path)
    elif format == "json":
        return save_json(presentation_data, output_path)
    else:
        logger.error(f"Unsupported export format: {format}")
        return False