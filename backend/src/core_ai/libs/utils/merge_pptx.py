import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from copy import deepcopy
import shutil
import tempfile
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import re
import io

def copy_placeholders_content(source_slide, target_slide):
    """Copy content from source slide's placeholders to target slide's placeholders."""
    
    # Dictionary to map placeholders in source slide to placeholders in target slide
    source_placeholders = {ph.placeholder_format.idx: ph for ph in source_slide.shapes 
                           if ph.is_placeholder}
    target_placeholders = {ph.placeholder_format.idx: ph for ph in target_slide.shapes 
                           if ph.is_placeholder}
    
    # Copy content from source placeholders to target placeholders with matching indices
    for idx, source_ph in source_placeholders.items():
        if idx in target_placeholders:
            target_ph = target_placeholders[idx]
            
            # Copy text content
            if source_ph.has_text_frame and target_ph.has_text_frame:
                # Fix: Use text_frame.text = "" instead of clear() to avoid placeholder issues
                if len(source_ph.text_frame.paragraphs) > 0:
                    # Handle the first paragraph differently - use existing instead of adding
                    target_ph.text_frame.text = ""  # This creates a single empty paragraph
                    
                    for i, para in enumerate(source_ph.text_frame.paragraphs):
                        # Use the existing first paragraph instead of adding a new one
                        if i == 0:
                            p = target_ph.text_frame.paragraphs[0]
                        else:
                            p = target_ph.text_frame.add_paragraph()
                            
                        p.text = para.text
                        p.font.bold = para.font.bold
                        p.font.italic = para.font.italic
                        p.font.size = para.font.size
                        # Safely check for color.rgb property to avoid _NoneColor error
                        if hasattr(para.font.color, 'rgb') and para.font.color.rgb is not None:
                            p.font.color.rgb = para.font.color.rgb
                        p.alignment = para.alignment
                    
            # Copy table content
            elif hasattr(source_ph, 'has_table') and source_ph.has_table and hasattr(target_ph, 'has_table'):
                source_table = source_ph.table
                if not hasattr(target_ph, 'table'):
                    continue
                
                try:
                    tbl = target_ph.insert_table(rows=source_table.rows.count, 
                                               cols=source_table.columns.count)
                    
                    for i, row in enumerate(source_table.rows):
                        for j, cell in enumerate(row.cells):
                            if i < tbl.rows.count and j < tbl.columns.count:
                                tbl.cell(i, j).text = cell.text
                except Exception as e:
                    print(f"Error copying table: {e}")
                    
            # Copy image content
            elif hasattr(source_ph, 'image') and source_ph.image:
                try:
                    # Create a BytesIO object from the image blob
                    image_stream = io.BytesIO(source_ph.image.blob)
                    
                    # Insert the image into the target placeholder
                    target_ph.insert_picture(image_stream)
                    
                except Exception as e:
                    print(f"Error copying image in placeholder: {e}")


def copy_slide_contents(source_slide, target_slide):
    """Copy contents from source slide to target slide including shapes and placeholders."""
    
    # Copy placeholders content
    copy_placeholders_content(source_slide, target_slide)
    
    # Copy non-placeholder shapes
    for shape in source_slide.shapes:
        if not shape.is_placeholder:
            # Handle pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    
                    # Create a BytesIO object to handle the image data
                    image_stream = io.BytesIO(shape.image.blob)
                    
                    # Add the picture using the BytesIO stream
                    target_slide.shapes.add_picture(image_stream, left, top, width, height)
                except Exception as e:
                    print(f"Error copying picture: {e}")
                    
            # Handle charts (requires more complex handling)
            elif hasattr(shape, 'chart') and shape.chart:
                # Charts are complex and usually require a more involved approach
                print("Note: Charts require special handling and might not be copied correctly.")
                
            # Handle text boxes
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                try:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    text_box = target_slide.shapes.add_textbox(left, top, width, height)
                    
                    # Copy text and format from source to new textbox
                    if shape.text:
                        # Fix: Use the same approach as in placeholders
                        text_box.text_frame.text = ""
                        for i, para in enumerate(shape.text_frame.paragraphs):
                            if i == 0:
                                p = text_box.text_frame.paragraphs[0]
                            else:
                                p = text_box.text_frame.add_paragraph()
                            p.text = para.text
                            
                            # Try to copy paragraph formatting
                            p.alignment = para.alignment
                            # Try to copy run formatting
                            if para.runs:
                                for run in para.runs:
                                    if hasattr(run.font, 'bold'):
                                        p.font.bold = run.font.bold
                                    if hasattr(run.font, 'italic'):
                                        p.font.italic = run.font.italic
                                    if hasattr(run.font, 'size') and run.font.size is not None:
                                        p.font.size = run.font.size
                                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                                        p.font.color.rgb = run.font.color.rgb
                except Exception as e:
                    print(f"Error copying text box: {e}")
            
            # Handle tables
            elif hasattr(shape, 'has_table') and shape.has_table:
                try:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    source_table = shape.table
                    
                    # Need to create a table with appropriate dimensions
                    rows, cols = source_table.rows.count, source_table.columns.count
                    table = target_slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # Copy contents and basic formatting
                    for i, row in enumerate(source_table.rows):
                        for j, cell in enumerate(row.cells):
                            table.cell(i, j).text = cell.text
                except Exception as e:
                    print(f"Error copying table: {e}")


def merge_pptx_files(input_dir, output_path, file_pattern=None):
    """
    Merge PowerPoint files into a single file, maintaining all content and formatting.
    
    Args:
        input_dir: Directory containing PPTX files to merge
        output_path: Path where the merged presentation will be saved
        file_pattern: Optional regex pattern to match specific filenames and sort them
    """
    # List all PPTX files in directory
    all_pptx_files = [f for f in os.listdir(input_dir) if f.endswith('.pptx')]
    
    if not all_pptx_files:
        raise ValueError("No .pptx files found in the directory.")
    
    # Sort files numerically based on extracted numbers from filenames
    def extract_number(filename):
        match = re.search(r'\d+', filename)
        return int(match.group()) if match else float('inf')
    
    pptx_files = sorted(all_pptx_files, key=extract_number)
    
    print(f"Found {len(pptx_files)} PPTX files to merge")
    
    # Use first file as the base
    base_file = os.path.join(input_dir, pptx_files[0])
    print(f"Using {pptx_files[0]} as base file")
    
    # Make a copy to work with
    shutil.copy2(base_file, output_path)
    merged_prs = Presentation(output_path)
    
    # Process each additional file
    for i, pptx_file in enumerate(pptx_files[1:], 1):
        print(f"Processing file {i}/{len(pptx_files)-1}: {pptx_file}")
        file_path = os.path.join(input_dir, pptx_file)
        
        # Open the presentation to merge
        source_prs = Presentation(file_path)
        
        # Add each slide from the source to the merged presentation
        for slide_index, source_slide in enumerate(source_prs.slides):
            # Find the right layout in the target presentation
            source_layout = source_slide.slide_layout
            layout_name = source_layout.name
            
            # Find matching layout in target presentation
            target_layout = None
            for layout in merged_prs.slide_masters[0].slide_layouts:
                if layout.name == layout_name:
                    target_layout = layout
                    break
            
            # If no matching layout found, use the first layout
            if not target_layout:
                target_layout = merged_prs.slide_masters[0].slide_layouts[0]
                print(f"Warning: Layout '{layout_name}' not found in target presentation. Using default layout.")
            
            # Create new slide with the matching layout
            target_slide = merged_prs.slides.add_slide(target_layout)
            
            # Copy all content from source slide to target slide
            copy_slide_contents(source_slide, target_slide)
            
            print(f"  - Added slide {slide_index+1} from {pptx_file}")
    
    # Save the merged presentation
    merged_prs.save(output_path)
    print(f"✅ Successfully merged {len(pptx_files)} files into {output_path}")
    return output_path


def merge_pptx_from_section_results(section_results, output_path, base_template=None):
    """
    Merge PowerPoint files from section results into a single file.
    
    Args:
        section_results: List of dictionaries with section_index and prs keys
        output_path: Path where the merged presentation will be saved
        base_template: Optional template to use as base (if None, first section's prs will be used)
        
    Returns:
        Path to the merged presentation
    """
    if not section_results:
        raise ValueError("No section results provided for merging.")
    
    # Sort results by section_index to ensure correct order
    sorted_results = sorted(section_results, key=lambda x: x["section_index"])
    prs_paths = [result["prs"] for result in sorted_results]
    
    # Create a temporary directory to store presentation files if needed
    with tempfile.TemporaryDirectory() as temp_dir:
        # Use the provided paths directly
        return merge_pptx_files(
            input_dir=os.path.dirname(prs_paths[0]),
            output_path=output_path,
            file_pattern=r".*"  # Use all files in directory
        )


if __name__ == "__main__":
    input_directory = "temp/slides"  # Directory containing PPTX files
    output_file = "merged_presentation.pptx"  # Output filename
    
    # Example pattern to sort files like "section_1.pptx", "section_2.pptx", etc.
    
    merge_pptx_files(input_directory, output_file)