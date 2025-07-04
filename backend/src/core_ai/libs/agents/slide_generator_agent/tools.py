from typing import Dict, Any, List, Optional, Union
from pydantic import BaseModel, Field
from langchain.tools import BaseTool
from pptx import Presentation
import os
import json
import logging

import re
from .schemas import PresentationContent, SlideContent, SlideLayoutManager
from .create_presetation import create_presentation
from core_ai.libs.utils.pptx2pdf import convert_to_pdf

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger("MasterSlideTool")


class MasterSlideTool(BaseTool):
    """Tool for matching PowerPoint templates and placeholders appropriately"""
    
    template_path : str = Field(description="Path to PowerPoint template file")
    name : str = "master_slide_tool"
    description : str  = """
    This tool helps create PowerPoint presentations by matching templates and placeholders appropriately.
    It can analyze templates, suggest layouts, and generate slides based on content.
    """

    def _run(self):
        pass
    
    def get_formatted_layouts_for_prompting(self, indexes: List[int] = None) -> str:
        """Get a formatted string of selected layouts for prompting"""
        result = f"INFORMATION ABOUT LAYOUTS IN TEMPLATE {os.path.basename(self.template_path)}:\n\n"
        
        try:
            presentation = Presentation(self.template_path)
            
            for master_idx, master in enumerate(presentation.slide_masters):
                for layout_idx, layout in enumerate(master.slide_layouts):
                    if indexes is not None and layout_idx not in indexes:
                        continue
                    
                    result += f"Layout {layout_idx} - {layout.name}:\n"
                    result += f"  Contains {len(layout.placeholders)} placeholder(s)\n"
                    
                    for ph in layout.placeholders:
                        ph_type = str(ph.placeholder_format.type)
                        result += f"  * {ph.name} (Type: {ph_type})\n"
                    
                    result += "\n"
            
            logger.info(f"Generated formatted layouts information for {self.template_path}")
            return result
            
        except Exception as e:
            logger.error(f"Error generating formatted layouts: {str(e)}")
            return f"Error getting layout information: {str(e)}"
            
    def create_map(self) -> Dict[str, Dict[str, int]]:
        """
        Create a mapping of layouts to their placeholders and type numbers.
        Handles special characters in layout names properly.
        
        Returns:
            Dict[str, Dict[str, int]]: Dictionary with layout names as keys and 
                                    dictionaries of placeholder names to type numbers as values
        """
        try:
            layout_map = {}
            
            # Get all layouts info
            presentation = Presentation(self.template_path)
            
            for master in presentation.slide_masters:
                for layout in master.slide_layouts:
                    placeholder_map = {}
                    
                    for placeholder in layout.placeholders:
                        key = placeholder.name
                        ph_type = str(placeholder.placeholder_format.type)
                        
                        # Extract the type number from the string format
                        try:
                            # Assuming the format is something like "PP_TYPE (12)"
                            type_num = int(ph_type.split('(')[-1].split(')')[0])
                        except (IndexError, ValueError):
                            # Fallback if the format is different
                            type_num = int(ph_type[-2:-1]) if ph_type[-2:-1].isdigit() else 0
                        
                        placeholder_map[key] = type_num
                    
                    # Use layout name as key, keeping special characters intact
                    layout_name = layout.name
                    layout_map[layout_name] = placeholder_map
            
            logger.info(f"Created layout to placeholder type mapping for {self.template_path}")
            return layout_map
            
        except Exception as e:
            logger.error(f"Error creating map: {str(e)}")
            return {}


        
        except Exception as e:
            logger.error(f"Error in mapping_index: {str(e)}")
            return None

    


        

        
        