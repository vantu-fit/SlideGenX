from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from pptx.enum.shapes import PP_PLACEHOLDER
import math
import os
import io
from typing import List, Dict, Optional, Any
import re 


class SlideHandler:
    def __init__(self, template_path: str, normal_font_size_pt: int = 18, title_font_size_pt: int = 24):
        self.prs = Presentation(template_path)
        self.template_path = template_path
        self.normal_font_size = normal_font_size_pt
        self.title_font_size = title_font_size_pt

    def master_layout_info(self, layout, index: int):
        result = f"Layout Index {index}: {layout.name}\n"
        for shape in layout.shapes:
            if not shape.is_placeholder:
                continue
            phf = shape.placeholder_format
            result += f"  - Placeholder index: {phf.idx}, type: {phf.type}, name: {shape.name}\n"
        return result

    def to_llm(self, indexs: List[int] = None):
        results = "INFORMATION ABOUT LAYOUTS:\n"
        # Nếu indexs không được truyền hoặc rỗng, lấy tất cả layouts
        if indexs is None or len(indexs) == 0:
            indexs = range(len(self.prs.slide_master.slide_layouts))

        for index in indexs:
            try:
                layout = self.prs.slide_master.slide_layouts[index]
                results += self.master_layout_info(layout, index)
            except IndexError:
                results += f"Error: Layout index {index} does not exist.\n"

        return results

    def to_llm_json(self, indexs: List[int]) -> Dict:
        results = []
        # Nếu indexs không được truyền hoặc rỗng, lấy tất cả layouts
        if indexs is None or len(indexs) == 0:
            indexs = range(len(self.prs.slide_master.slide_layouts))

        for index in indexs:
            try:
                item = {}
                layout = self.prs.slide_master.slide_layouts[index]
                for shape in layout.shapes:
                    if not shape.is_placeholder:
                        continue
                    phf = shape.placeholder_format
                    index = f"placeholder index {phf.idx}"
                    type = phf.type
                    name = shape.name
                    item[index] = {"type": type, "name": name}
                results.append(item)

            except IndexError:
                results += f"Error: Layout index {index} does not exist.\n"

        return results


    def _get_font_size_for_placeholder(self, placeholder):
        """
        Determine the appropriate font size based on the placeholder name
        If the placeholder name contains 'title' (case insensitive), use title font size,
        otherwise use normal font size
        """
        if placeholder.name and 'title' in placeholder.name.lower():
            return self.title_font_size
        return self.normal_font_size

    def _estimate_text_fit(self, placeholder, total_char_count):
        width = placeholder.width / 12700
        height = placeholder.height / 12700

        font_size = self._get_font_size_for_placeholder(placeholder)
        char_width_pt = font_size * 0.5
        line_height_pt = font_size * 1.3
        chars_per_line = max(int(width / char_width_pt), 1)
        max_lines = max(int(height / line_height_pt),1)
        max_chars = chars_per_line * max_lines

        overflow = total_char_count > max_chars
        return {
            "overflow": overflow,
            "max_lines": max_lines,
            "chars_per_line": chars_per_line,
            "max_chars": max_chars,
        }

    def _estimate_list_text_fit(self, placeholder, list_str: list[str]):
        """
        Ước lượng số dòng và số ký tự mỗi dòng cho mỗi item trong danh sách.
        :param placeholder: Placeholder của slide
        :param list_str: Danh sách các chuỗi văn bản
        :return: List các thông tin ước lượng cho từng item (số dòng và số ký tự mỗi dòng)
        """
        lines_required = 0
        width = placeholder.width / 12700
        height = placeholder.height / 12700

        font_size = self._get_font_size_for_placeholder(placeholder)
        char_width_pt = font_size * 0.5  # approximate width per char
        line_height_pt = font_size * 1.5  # line spacing
        chars_per_line = max(int(width / char_width_pt), 1)
        max_lines = int(height / line_height_pt)
        max_chars = chars_per_line * max_lines

        for text in list_str:
            lines_required += max(math.ceil(len(text) / chars_per_line), 1)

        overflow = lines_required > max_lines

        return {
            "overflow": overflow,
            "max_lines": max_lines,
            "chars_per_line": chars_per_line,
            "max_chars": max_chars,
            "type": "list",
        }

    def insert_text_str(self, placeholder, text: str):
        estimation = self._estimate_text_fit(placeholder, len(text))
        if estimation["overflow"]:
            expected_text = text[: estimation["max_chars"]]
            return {
                "fit": False,
                "max_lines": estimation["max_lines"],
                "chars_per_line": estimation["chars_per_line"],
                "max_chars": estimation["max_chars"],
                "type": "text",
                "text": text,
            }
            
        # Insert text and set the font size
        placeholder.text = text
        
        # Apply font size to all text runs in the placeholder
        font_size = self._get_font_size_for_placeholder(placeholder)
        for paragraph in placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                
        return {"fit": True}

    def insert_text_list(self, placeholder, texts: list[str]):
        max_info = self._estimate_list_text_fit(placeholder, texts)

        if max_info["overflow"]:
            fit_text = []
            total_lines = 0
            for line in texts:
                line_len = len(line)
                line_max_chars = min(line_len, max_info["chars_per_line"])
                fit_text.append(line[:line_max_chars])
                total_lines += 1
                if total_lines >= max_info["max_lines"]:
                    break

            return {
                "fit": False,
                "max_lines": max_info["max_lines"],
                "chars_per_line": max_info["chars_per_line"],
                "max_chars": max_info["max_chars"],
                "type": "list",
                "text": texts,
            }

        # Insert the text
        placeholder.text = "\n".join(texts)
        
        # Apply font size to all text runs in the placeholder
        font_size = self._get_font_size_for_placeholder(placeholder)
        for paragraph in placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                
        return {"fit": True}

    def _add_padding_to_image(self, image_path, target_ratio):
        """
        Add padding to an image to match the target aspect ratio without cropping.
        
        Args:
            image_path: Path to the original image
            target_ratio: The desired width/height ratio
            
        Returns:
            BytesIO object with the padded image
        """
        with Image.open(image_path) as img:
            orig_width, orig_height = img.size
            orig_ratio = orig_width / orig_height
            
            # Determine if we need horizontal or vertical padding
            if orig_ratio < target_ratio:
                # Need to add horizontal padding
                new_width = int(orig_height * target_ratio)
                new_height = orig_height
                
                # Create a new image with white background
                new_img = Image.new("RGB", (new_width, new_height), (255, 255, 255))
                
                # Calculate position to paste original image (centered)
                paste_x = (new_width - orig_width) // 2
                paste_y = 0
                
                # Paste original image onto the new image
                new_img.paste(img, (paste_x, paste_y))
                
            elif orig_ratio > target_ratio:
                # Need to add vertical padding
                new_width = orig_width
                new_height = int(orig_width / target_ratio)
                
                # Create a new image with white background
                new_img = Image.new("RGB", (new_width, new_height), (255, 255, 255))
                
                # Calculate position to paste original image (centered)
                paste_x = 0
                paste_y = (new_height - orig_height) // 2
                
                # Paste original image onto the new image
                new_img.paste(img, (paste_x, paste_y))
                
            else:
                # No padding needed, just return the original image
                new_img = img.copy()
            
            # Save the padded image to a BytesIO object
            output = io.BytesIO()
            new_img.save(output, format=img.format if img.format else "PNG")
            output.seek(0)
            
            return output

    def insert_image(self, placeholder, image_path: str, threshold=0.3):
        try:
            # Initialize variables for error reporting
            ph_ratio = None
            img_ratio = None
            
            # Check if the image exists
            if not os.path.exists(image_path):
                raise FileNotFoundError(f"Image not found: {image_path}")
            
            # Get image dimensions and ratio
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                img_ratio = img_width / img_height

            # Validate placeholder
            if (
                placeholder.shape_type != MSO_SHAPE_TYPE.PICTURE
                and not placeholder.placeholder_format
            ):
                raise ValueError("Invalid placeholder for image")

            # Get placeholder dimensions and ratio
            ph_width = placeholder.width
            ph_height = placeholder.height
            ph_ratio = ph_width / ph_height

            # Adjust threshold for extreme aspect ratios
            if ph_ratio >= 3:
                threshold *= 6

            # If the difference is small enough, just insert the image directly
                # Add padding to match the placeholder's aspect ratio
            padded_image = self._add_padding_to_image(image_path, ph_ratio)
            
            # Insert the padded image into the placeholder
            placeholder.insert_picture(padded_image)
            
            return {"fit": True}
                
        except Exception as e:
            # Error handling - make sure we have values for ph_ratio and img_ratio
            error_info = {
                "fit": False,
                "type": "image",
                "image_destination": image_path,
                "error": str(e)
            }
            
            # Only add ratios if they were successfully calculated
            if ph_ratio is not None:
                error_info["expected_image_ratio"] = round(ph_ratio, 2)
            if img_ratio is not None:
                error_info["actual_image_ratio"] = round(img_ratio, 2)
                
            return error_info

    def get_type_slide(self, layout_index, placeholder_index):
        layout = self.prs.slide_master.slide_layouts[layout_index]
        for shape in layout.shapes:
            if not shape.is_placeholder:
                continue
            phf = shape.placeholder_format
            if phf.idx == placeholder_index:
                return phf.type
        return None


    def create_slides(
        self,
        slides: List[Dict[str, Any]],
        return_prs: bool = False,
        output_path: Optional[str] = "output.pptx",
    ) -> List[Dict[str, Any]]:
        prs = Presentation(self.template_path)
        results = []
        images = []
        for slide_index, slide in enumerate(slides):
            list_results = []
            layout_index = slide.get("layout_index", 0)
            layout = prs.slide_master.slide_layouts[layout_index]

            slide_obj = prs.slides.add_slide(layout)
            is_insert_image = False

            for placeholder_idx, content in slide.items():
                if placeholder_idx == "layout_index":
                    continue

                if isinstance(placeholder_idx, str):
                    match = re.search(r'\d+', placeholder_idx)
                    if match:
                        placeholder_idx = int(match.group())

                place_folder_type = self.get_type_slide(
                    layout_index, int(placeholder_idx)
                )
                if place_folder_type:
                    place_folder_type = place_folder_type.name
                if place_folder_type not in [
                    "CENTER_TITLE",
                    "TITLE",
                    "SUBTITLE",
                    "TITLE_AND_CONTENT",
                    "TITLE_AND_OBJECT",
                    "BODY",
                    "OBJECT",
                    "PICTURE",
                    "DIAGRAM",
                    "CHART",
                ]:
                    continue

                placeholder = slide_obj.placeholders[int(placeholder_idx)]

                if isinstance(content, str):
                    if content.endswith(".png") or content.endswith(".jpg"):
                        insert_info = self.insert_image(placeholder, content)
                        insert_info["placeholder"] = placeholder_idx
                        insert_info["slide_index"] = slide_index
                        insert_info["type"] = "image" if "image" in content else "diagram"
                        insert_info["path"] = content
                        is_insert_image = insert_info
                    else:
                        insert_info = self.insert_text_str(placeholder, content)
                        insert_info["placeholder"] = placeholder_idx
                        list_results.append(insert_info)
                elif isinstance(content, list):
                    insert_info = self.insert_text_list(placeholder, content)
                    insert_info["placeholder"] = placeholder_idx
                    list_results.append(insert_info)
                    images.append({"fit": True})
            if is_insert_image:
                images.append(is_insert_image)

            results.append(list_results)

        def filter_func(result):
            return [r for r in result if r["fit"] is False]

        if return_prs:
            prs.save(output_path)
            print(f"Saved presentation to {output_path}")
            return [filter_func(result) for result in results], [r if r["fit"] is False else [] for r in images] , output_path
        else:
            return [filter_func(result) for result in results], [r if r["fit"] is False else [] for r in images], prs

    def resutl_to_llm(self, results: List[Dict[str, Any]]) -> str:
        result = "SUMMARY OF SLIDE CREATION:\n"
        for i, slide in enumerate(results):
            result += f"""
            Slide {i + 1}:
            """
            for placeholder in slide:
                if placeholder["type"] == "list":
                    requirements = f"""
                        placeholder: {placeholder['placeholder']}
                            - Content: {placeholder['text']}
                            - Requirement: "Convert Content Shoter" Summarize content in to max lines is {placeholder['max_lines']} and max chars per line is {placeholder['chars_per_line']}
                    """
                elif placeholder["type"] == "text":
                    requirements = f"""
                        placeholder: {placeholder['placeholder']}
                            - Content: {placeholder['text']}
                            - Requirement: "Convert Content Shoter" Summarize content in to max chars is {placeholder['max_chars']}
                    """
                else:
                    continue
                result += requirements
        result += "\nIF SLIDE THAT NOT HAVE CONTENT THEN JUST RETURN SLIDE NUMBER\n\n"

        return result
    
    def result_visual_to_llm(self, results: List[Dict[str, Any]]) -> List[str]:
        prompts = []
        for i, insert_info in enumerate(results):
            if insert_info == [] or "type" not in insert_info:
                prompts.append("")
                continue
            if insert_info["type"] == "diagram":
                prompt = f"""IMPORTANT: Change to another type of diagram or redraw the diagram for width divide hight ratio is {insert_info["expected_image_ratio"]} and actual image ratio is {insert_info["actual_image_ratio"]}"""
                prompts.append(prompt)
            elif insert_info["type"] == "image":
                prompt = f"""IMPORTANT: Find another image for width and hight ratio is {insert_info["expected_image_ratio"]} and actual image ratio is {insert_info["actual_image_ratio"]}"""
                prompts.append(prompt)
            else:
                prompts.append("")
        return prompts


if __name__ == "__main__":

    # Khởi tạo handler với font size mặc định (18pt cho normal text, 24pt cho title)
    handler = SlideHandler("pptx_templates/Minimalist.pptx")

    print(handler.to_llm())
