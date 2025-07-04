import os, datetime
from pptx import Presentation
from .schemas import PresentationContent, SlideLayoutManager
from libs.utils.pptx2pdf import convert_to_pdf
import json


def create_presentation(template_path: str, slide_plan: PresentationContent) -> Presentation:
    """Create presentation using the slide plan with comprehensive placeholder mapping and cleanup of empty placeholders"""
    prs = Presentation(template_path)
    
    print(f"Creating presentation with {len(slide_plan.slides)} slides")
    layout_manager = SlideLayoutManager(template_path)
    # Create a placeholder mapping for each layout
    layout_placeholder_map = layout_manager.create_map()

    print(f"{json.dumps(layout_placeholder_map, indent=2)}")

    placeholder_type_map = {
        "TITLE": 1,
        "BODY": 2,
        "CENTER_TITLE": 3,
        "OBJECT": 7,
        "PICTURE": 18,
        "SLIDE_NUMBER": 13
    }    
    
    # Dictionary mapping to translate placeholder types from strings to numbers
    
    for slide_index, slide_content in enumerate(slide_plan.slides):
        # Get the layout by index (adjusted for 0-based indexing)
        layout_idx = slide_content.layout_index - 1
        if layout_idx < 0 or layout_idx >= len(prs.slide_layouts):
            print(f"Warning: Invalid layout index {slide_content.layout_index} for slide {slide_index+1}. Using default layout.")
            layout_idx = 0  # Use the first layout as default
            
        slide_layout = prs.slide_layouts[layout_idx]
        print(f"\nCreating slide {slide_index+1} with layout index {slide_content.layout_index} ({layout_idx})")
        
        # Add a new slide with the selected layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Get placeholder mapping for this layout
        layout_map = layout_placeholder_map[slide_content.layout_index - 1]
        
        # Log actual placeholders in slide for debugging
        actual_placeholders = []
        for shape in slide.placeholders:
            if hasattr(shape, "placeholder_format"):
                ph_name = shape.name if hasattr(shape, "name") else "Unnamed"
                ph_idx = shape.placeholder_format.idx
                ph_type = shape.placeholder_format.type
                actual_placeholders.append((ph_name, ph_idx, ph_type))
                print(f"Actual placeholder in slide: {ph_name} (idx: {ph_idx}, type: {ph_type})")
        
        print(f"Available content keys: {list(slide_content.content.keys())}")
        
        # Create a mapping from layout placeholder names to actual placeholder indices
        placeholder_map = {}
        
        # First, try to map by type
        for content_key, type_name in layout_map.items():
                
            # Find matching placeholder by type
            for ph_name, ph_idx, ph_type in actual_placeholders:
                if ph_type == type_name:
                    placeholder_map[content_key] = (ph_name, ph_idx, ph_type)
                    # delete actual_placeholders.remove(placeholder_info)
                    # actual_placeholders.remove((ph_name, ph_idx, ph_type))
                    break
        
        # Special handling for Content Placeholder and OBJECT type
        # This is needed because there might be multiple OBJECT type placeholders
        object_placeholders = [(ph_name, ph_idx, ph_type) 
                            for ph_name, ph_idx, ph_type in actual_placeholders 
                            if ph_type == placeholder_type_map.get("OBJECT")]
        
        # Map Content Placeholder 1, 2, 3, etc. to actual object placeholders
        object_count = 1
        for ph_name, ph_idx, ph_type in object_placeholders:
            content_key = f"Content Placeholder {object_count}"
            if content_key in slide_content.content and content_key not in placeholder_map:
                placeholder_map[content_key] = (ph_name, ph_idx, ph_type)
                object_count += 1
        
        # Handle title specially
        title_placeholders = [(ph_name, ph_idx, ph_type) 
                            for ph_name, ph_idx, ph_type in actual_placeholders 
                            if ph_type in [placeholder_type_map.get("TITLE"), 
                                        placeholder_type_map.get("CENTER_TITLE")]]
        
        if title_placeholders and "Title 1" in slide_content.content:
            placeholder_map["Title 1"] = title_placeholders[0]
        
        # Print the placeholder map for debugging
        print(f"Placeholder map: {placeholder_map}")
        
        # Track which placeholders have been filled
        filled_placeholders = set()
        
        # Now, fill in the placeholders
        for content_key, content in slide_content.content.items():
            if content is None or content == "" or content == " ":
                continue
                
            # Find matching placeholder
            placeholder_info = placeholder_map.get(content_key)
            
            if placeholder_info:
                ph_name, ph_idx, ph_type = placeholder_info
                
                # Find the shape
                shape = None
                for s in slide.shapes:
                    if (hasattr(s, "placeholder_format") and 
                        s.placeholder_format.idx == ph_idx):
                        shape = s
                        break
                
                if not shape:
                    print(f"Could not find shape for placeholder: {ph_name}")
                    continue
                
                print(f"Filling placeholder {ph_name} (idx: {ph_idx}) with content from {content_key}")
                
                # Mark this placeholder as filled
                filled_placeholders.add(ph_idx)
                
                # Fill placeholder based on type
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    text_frame.clear()  # Clear any default text
                    
                    if isinstance(content, list):
                        # Handle bullet points
                        for i, item in enumerate(content):
                            # Handle nested lists with dictionaries
                            if isinstance(item, dict) and "text" in item and "level" in item:
                                if i == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                    
                                p.text = item["text"]
                                p.level = item["level"]
                            # Handle simple list items
                            else:
                                if i == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                    
                                p.text = str(item)
                                # Set bullet formatting for body placeholders
                                if ph_type in [2, 7]:  # BODY or OBJECT type
                                    p.level = 0
                    else:
                        # Plain text
                        text_frame.paragraphs[0].text = str(content)
                        
                        # For titles, adjust font size and alignment
                        if ph_type in [1, 3]:  # TITLE or CENTER_TITLE
                            p = text_frame.paragraphs[0]
                            p.alignment = 1  # Center alignment
                            
                            # Bold title text
                            for run in p.runs:
                                run.font.bold = True
                    
                    print(f"Added text content to placeholder {ph_name}")
                
                elif hasattr(shape, "table") and isinstance(content, list) and all(isinstance(row, list) for row in content):
                    # Handle table content
                    table = shape.table
                    # ... [table handling code remains the same]
                
                elif ph_type == 18 and isinstance(content, str):  # PICTURE type
                    # Handle picture content
                    if os.path.exists(content):
                        try:
                            shape.insert_picture(content)
                            print(f"Inserted picture from {content} into placeholder {ph_name}")
                        except Exception as e:
                            print(f"Error inserting picture: {e}")
                    else:
                        print(f"Picture content is not a valid file path: {content}")
                
            else:
                print(f"No placeholder found for content key: {content_key}")
        
        # Clear or hide empty placeholders
        for shape in slide.shapes:
            if (hasattr(shape, "placeholder_format") and 
                shape.placeholder_format.idx not in filled_placeholders and
                shape.placeholder_format.type != placeholder_type_map.get("SLIDE_NUMBER")):  # Preserve slide numbers
                
                ph_name = shape.name if hasattr(shape, "name") else "Unnamed"
                ph_idx = shape.placeholder_format.idx
                ph_type = shape.placeholder_format.type
                
                print(f"Cleaning empty placeholder: {ph_name} (idx: {ph_idx}, type: {ph_type})")
                
                try:
                    # For text placeholders, clear the text
                    if hasattr(shape, "text_frame"):
                        # Check if it has default text (often the case with placeholders)
                        has_text = False
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                has_text = True
                                break
                        
                        if has_text:
                            shape.text_frame.clear()
                            print(f"Cleared text from placeholder {ph_name}")
                    
                    # Try to make the shape invisible
                    # Note: Not all shapes can have transparency set, so we use a try-except
                    try:
                        if hasattr(shape, "fill"):
                            shape.fill.transparency = 1.0  # Make fully transparent
                        if hasattr(shape, "line"):
                            shape.line.transparency = 1.0  # Make border transparent
                    except Exception as e:
                        print(f"Could not make shape transparent: {e}")
                
                except Exception as e:
                    print(f"Error cleaning placeholder: {e}")
                
        # Add notes if available
        if hasattr(slide_content, "notes") and slide_content.notes:
            if not slide.has_notes_slide:
                slide.notes_slide
            slide.notes_slide.notes_text_frame.text = slide_content.notes
    
    # Save metadata
    if slide_plan.title:
        prs.core_properties.title = slide_plan.title
    else:
        prs.core_properties.title = "Generated Presentation"
    
    prs.core_properties.comments = f"Created with SlideGenerateAgent on {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    return prs

if __name__ == "__main__":
    # Example usage
    template_path = "pptx_templates/Bosch-WeeklyReport.pptx"
    slide_plan = PresentationContent.from_json_file("draft.json")
    
    prs = create_presentation(template_path, slide_plan)

    prs.save("output.pptx")